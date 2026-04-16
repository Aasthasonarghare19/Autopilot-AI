import streamlit as st
from pptx import Presentation

# -------- AI AGENT FUNCTION --------
def academic_agent(topic):
    plan = f"""
1. Introduction to {topic}
2. Key concepts
3. Applications
4. Advantages
5. Future scope
"""

    report = f"""
{topic} is a very important and rapidly growing field in modern technology. It focuses on improving efficiency, automation, and intelligent decision-making in different industries.

The main idea behind {topic} is to solve real-world problems using advanced techniques and systems. It is widely used in sectors such as healthcare, education, business, and engineering.

One of the key features of {topic} is its ability to process large amounts of data and generate useful insights. This helps organizations make better and faster decisions.

There are many advantages of {topic}, such as increased productivity, reduced human effort, and improved accuracy. However, it also has some challenges like high cost and complexity.

In the future, {topic} will play a major role in innovation and technological development. Many companies and researchers are investing heavily in this field to build smarter systems.
"""

    ppt_points = [
        f"Introduction to {topic}",
        "Definition and Overview",
        "Key Concepts",
        "Applications",
        "Advantages",
        "Challenges",
        "Future Scope",
        "Conclusion"
    ]

    summary = f"{topic} is a powerful and growing field with many real-world applications."

    return plan, report, ppt_points, summary


# -------- PPT GENERATOR --------
def create_ppt(topic, ppt_points):
    prs = Presentation()

    for point in ppt_points:
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)

        title = slide.shapes.title
        content = slide.placeholders[1]

        title.text = point
        content.text = f"""
- Explanation of {point}
- Key details related to {topic}
- Real-world applications
- Important points to understand
"""

    file_name = f"{topic}.pptx"
    prs.save(file_name)
    return file_name


# -------- UI --------
st.set_page_config(page_title="AutoPilot AI", layout="centered")

# -------- HEADER --------
st.markdown("""
<h1 style='text-align: center; color: #4CAF50;'>
🚀 AutoPilot AI – Academic Assistant
</h1>
<p style='text-align: center;'>
Generate Report, PPT & Summary Automatically
</p>
""", unsafe_allow_html=True)

# -------- INPUT --------
topic = st.text_input("📌 Enter your topic", placeholder="e.g. Artificial Intelligence")

# -------- BUTTON --------
if st.button("✨ Generate Content"):
    if topic:
        plan, report, ppt_points, summary = academic_agent(topic)

        # -------- PLAN --------
        st.markdown("## 📌 Plan")
        st.success(plan)

        # -------- REPORT --------
        st.markdown("## 📄 Detailed Report")
        st.info(report)

        # -------- PPT --------
        st.markdown("## 📊 PPT Content")
        st.write(ppt_points)

        # Generate PPT
        file = create_ppt(topic, ppt_points)

        # Download Button
        with open(file, "rb") as f:
            st.download_button(
                label="📥 Download PPT",
                data=f,
                file_name=file,
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )

        # -------- SUMMARY --------
        st.markdown("## 📌 Summary")
        st.warning(summary)

    else:
        st.error("⚠️ Please enter a topic")